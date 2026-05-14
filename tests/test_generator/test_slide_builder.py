"""Tests for SlideBuilder - Task 12."""

import pytest
from pptx import Presentation
from pptx.util import Inches

from html_pptx_template.templates.schema import Theme, LayoutSlide, LayoutElement
from html_pptx_template.generator.slide_builder import SlideBuilder


@pytest.fixture
def sample_theme():
    return Theme(
        colors={
            "primary": "#1E3A5F",
            "secondary": "#4A90D9",
            "accent": "#E67E22",
            "background": "#FFFFFF",
            "surface": "#F5F7FA",
            "text_primary": "#2C3E50",
            "text_secondary": "#7F8C8D",
            "text_on_primary": "#FFFFFF",
        },
        fonts={
            "heading": {
                "family": "Arial, sans-serif",
                "sizes": {"h1": 44, "h2": 32, "h3": 24},
                "weight": "bold",
                "color": "text_primary",
            },
            "body": {
                "family": "Arial, sans-serif",
                "size": 18,
                "weight": "normal",
                "color": "text_primary",
            },
            "caption": {
                "family": "Arial, sans-serif",
                "size": 14,
                "color": "text_secondary",
            },
        },
        spacing={
            "slide_padding": [40, 40, 40, 40],
            "content_gap": 20,
            "line_spacing": 1.5,
        },
        aspect_ratio="16:9",
        slide_width=10.0,
        slide_height=5.625,
    )


@pytest.fixture
def prs():
    return Presentation()


@pytest.fixture
def builder(sample_theme, prs):
    return SlideBuilder(sample_theme, prs)


class TestBuildTitleSlide:
    def test_build_title_slide(self, builder, prs):
        """Build title slide."""
        layout = LayoutSlide(
            name="title",
            slide_type="title",
            background="background",
            elements=[
                LayoutElement(type="title", style={"level": "h1", "align": "centered"}),
                LayoutElement(type="subtitle", style={"level": "body", "align": "centered"}),
            ],
        )
        content = {"title": "My Title", "subtitle": "My Subtitle"}
        slide = builder.build(layout, content)

        assert slide is not None
        # Should have shapes: background + title + subtitle
        assert len(slide.shapes) >= 2


class TestBuildContentSlide:
    def test_build_content_slide(self, builder, prs):
        """Build content slide."""
        layout = LayoutSlide(
            name="content",
            slide_type="content",
            background="background",
            elements=[
                LayoutElement(type="header_bar", style={"level": "h2", "align": "left"}),
                LayoutElement(type="content", style={"level": "body", "align": "left"}),
            ],
        )
        content = {"title": "Content Title", "body": "Some body text here."}
        slide = builder.build(layout, content)

        assert slide is not None
        assert len(slide.shapes) >= 2


class TestBuildWithBullets:
    def test_build_with_bullets(self, builder, prs):
        """Build bullet list slide."""
        layout = LayoutSlide(
            name="content",
            slide_type="content",
            background="background",
            elements=[
                LayoutElement(type="header_bar", style={"level": "h2", "align": "left"}),
                LayoutElement(type="content", style={"level": "body", "align": "left"}),
            ],
        )
        content = {
            "title": "Bullet Slide",
            "bullets": ["First point", "Second point", "Third point"],
        }
        slide = builder.build(layout, content)

        assert slide is not None
        assert len(slide.shapes) >= 2


class TestBuildSectionDivider:
    def test_build_section_divider(self, builder, prs):
        """Build section divider slide."""
        layout = LayoutSlide(
            name="section_divider",
            slide_type="section_divider",
            background="primary",
            elements=[
                LayoutElement(type="title", style={"level": "h1", "align": "centered"}),
            ],
        )
        content = {"title": "Section Title"}
        slide = builder.build(layout, content)

        assert slide is not None
        assert len(slide.shapes) >= 1


class TestBuildTwoColumn:
    def test_build_two_column(self, builder, prs):
        """Build two column slide."""
        layout = LayoutSlide(
            name="two_column",
            slide_type="two_column",
            background="background",
            elements=[
                LayoutElement(type="header_bar", style={"level": "h2", "align": "left"}),
                LayoutElement(type="column", style={"level": "body", "align": "left", "width_pct": "50%"}),
                LayoutElement(type="column", style={"level": "body", "align": "left", "width_pct": "50%"}),
            ],
        )
        content = {
            "title": "Two Column Title",
            "left": "Left column content.",
            "right": "Right column content.",
        }
        slide = builder.build(layout, content)

        assert slide is not None
        assert len(slide.shapes) >= 2
