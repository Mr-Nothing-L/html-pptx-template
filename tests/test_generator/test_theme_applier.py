"""Tests for ThemeApplier - Task 10."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from html_pptx_template.templates.schema import Theme
from html_pptx_template.generator.theme_applier import ThemeApplier


@pytest.fixture
def sample_theme():
    return Theme(
        colors={
            "primary": "#1E3A5F",
            "secondary": "#4A90D9",
            "accent": "#E67E22",
            "background": "#FFFFFF",
            "surface": "#F5F7FA",
            "text_primary": "#2C3E50",
            "text_secondary": "#7F8C8D",
            "text_on_primary": "#FFFFFF",
        },
        fonts={
            "heading": {
                "family": "Arial, sans-serif",
                "sizes": {"h1": 44, "h2": 32, "h3": 24},
                "weight": "bold",
                "color": "text_primary",
            },
            "body": {
                "family": "Arial, sans-serif",
                "size": 18,
                "weight": "normal",
                "color": "text_primary",
            },
            "caption": {
                "family": "Arial, sans-serif",
                "size": 14,
                "color": "text_secondary",
            },
        },
        spacing={
            "slide_padding": [40, 40, 40, 40],
            "content_gap": 20,
            "line_spacing": 1.5,
        },
        aspect_ratio="16:9",
        slide_width=10.0,
        slide_height=5.625,
    )


@pytest.fixture
def prs():
    return Presentation()


@pytest.fixture
def theme_applier(sample_theme):
    return ThemeApplier(sample_theme)


class TestApplySlideBackground:
    def test_apply_slide_background(self, theme_applier, prs):
        """Apply background color to slide."""
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        theme_applier.apply_slide_background(slide, color_key="background")

        background = slide.background
        fill = background.fill
        assert fill.type is not None


class TestApplyTextStyle:
    def test_apply_text_style_heading_h1(self, theme_applier, prs):
        """Apply heading text style, verify font/size/bold/color."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Test Heading"

        theme_applier.apply_text_style(run, text_type="heading", level="h1")

        assert run.font.size == Pt(44)
        assert run.font.bold is True
        assert run.font.name == "Arial"

    def test_apply_body_text(self, theme_applier, prs):
        """Body text has correct size and not bold."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Test Body"

        theme_applier.apply_text_style(run, text_type="body")

        assert run.font.size == Pt(18)
        assert run.font.bold is False
        assert run.font.name == "Arial"


class TestAddShapeWithFill:
    def test_add_shape_with_fill(self, theme_applier, prs):
        """Add a shape with solid fill from theme color."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        shape = theme_applier.add_shape_with_fill(
            slide, 1, Inches(1), Inches(1), Inches(2), Inches(0.5), color_key="primary"
        )
        assert shape is not None
        fill = shape.fill
        assert fill.type is not None


class TestGetSlideDimensions:
    def test_get_slide_dimensions(self, theme_applier):
        """Return (width, height) as Inches tuples."""
        width, height = theme_applier.get_slide_dimensions()
        assert width == Inches(10.0)
        assert height == Inches(5.625)
