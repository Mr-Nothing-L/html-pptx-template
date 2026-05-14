"""SlideBuilder - builds slides from layout + content."""

from pathlib import Path

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from html_pptx_template.generator.theme_applier import ThemeApplier


class SlideBuilder:
    """Build a slide from layout + content dict."""

    def __init__(self, theme, prs, assets_dir: Path = None):
        self.theme = theme
        self.prs = prs
        self.applier = ThemeApplier(theme)
        self.assets_dir = assets_dir
        self._width, self._height = self.applier.get_slide_dimensions()
        self._pad = theme.spacing.slide_padding  # [top, right, bottom, left]

    def build(self, layout, content: dict):
        """Build a slide from layout + content dict."""
        slide_layout = self.prs.slide_layouts[6]  # blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Apply background
        bg_color = layout.background or "background"
        self.applier.apply_slide_background(slide, color_key=bg_color)

        # Route to builder based on slide type
        slide_type = layout.slide_type
        if slide_type == "title":
            self._build_title_slide(slide, content)
        elif slide_type == "content":
            self._build_content_slide(slide, content)
        elif slide_type == "section_divider":
            self._build_section_divider(slide, content)
        elif slide_type == "two_column":
            self._build_two_column(slide, content)
        elif slide_type == "image_left":
            self._build_image_left(slide, content)
        elif slide_type == "image_full":
            self._build_image_full(slide, content)
        elif slide_type == "image_gallery":
            self._build_image_gallery(slide, content)
        elif slide_type == "table":
            self._build_table_slide(slide, content)
        elif slide_type == "chart":
            self._build_chart_slide(slide, content)
        else:
            self._build_content_slide(slide, content)

        return slide

    # --- Helper: common measurements ---

    def _content_rect(self):
        """Return (left, top, width, height) for content area inside padding."""
        left = Inches(self._pad[3] / 72)   # left pad in inches
        top = Inches(self._pad[0] / 72)    # top pad in inches
        right = self._width - Inches(self._pad[1] / 72)
        bottom = self._height - Inches(self._pad[2] / 72)
        return left, top, right - left, bottom - top

    def _resolve_image_path(self, path: str) -> Path | None:
        """Resolve an image path: absolute, relative to assets_dir, or relative to CWD."""
        p = Path(path)
        if p.is_absolute() and p.exists():
            return p
        if self.assets_dir:
            asset = self.assets_dir / p
            if asset.exists():
                return asset
        cwd = Path.cwd() / p
        if cwd.exists():
            return cwd
        return None

    def _add_textbox(self, slide, left, top, width, height, text: str,
                     text_type="body", level=None, align=PP_ALIGN.LEFT,
                     color_key=None, uppercase=False):
        """Add a textbox with theme styling."""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        self.applier.apply_text_style(run, text_type=text_type, level=level, uppercase=uppercase)
        if color_key:
            from html_pptx_template.utils.color import hex_to_rgb
            from pptx.dml.color import RGBColor
            r, g, b = hex_to_rgb(getattr(self.theme.colors, color_key))
            run.font.color.rgb = RGBColor(r, g, b)
        return box

    # --- Layout builders ---

    def _build_title_slide(self, slide, content: dict):
        """Editorial title: surface background, massive centered ALL CAPS title."""
        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        left, top, cw, ch = self._content_rect()

        # Background: surface color for editorial feel
        self.applier.apply_slide_background(slide, color_key="surface")

        # Title — centered vertically in upper half, massive, ALL CAPS
        title_h = Inches(1.8)
        title_top = top + Inches(0.8)
        self._add_textbox(
            slide, left, title_top, cw, title_h,
            title, text_type="heading", level="h1",
            align=PP_ALIGN.CENTER, uppercase=True,
        )

        # Subtitle — minimal, small, secondary color, centered below
        if subtitle:
            sub_h = Inches(0.6)
            sub_top = title_top + title_h + Inches(0.2)
            self._add_textbox(
                slide, left, sub_top, cw, sub_h,
                subtitle, text_type="caption",
                align=PP_ALIGN.CENTER, color_key="text_secondary",
            )

    def _build_content_slide(self, slide, content: dict):
        """Editorial content: light background, large heading, sparse bullets. No header bar."""
        title = content.get("title", "")
        body = content.get("body", "")
        bullets = content.get("bullets", [])

        left, top, cw, ch = self._content_rect()
        gap = Inches(self.theme.spacing.content_gap / 72)

        # No colored header bar — editorial minimal style
        current_top = top

        # Large left-aligned heading
        if title:
            title_h = Inches(0.9)
            self._add_textbox(
                slide, left, current_top, cw, title_h,
                title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        # Body / bullets with generous spacing
        content_h = ch - (current_top - top)
        if content_h > Inches(0.5):
            box = slide.shapes.add_textbox(left, current_top, cw, content_h)
            tf = box.text_frame
            tf.word_wrap = True

            if bullets:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(16)  # generous spacing
                    run = p.add_run()
                    run.text = bullet
                    self.applier.apply_text_style(run, text_type="body")
            elif body:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = body
                self.applier.apply_text_style(run, text_type="body")

    def _build_section_divider(self, slide, content: dict):
        """Dramatic divider: solid primary background, white centered ALL CAPS title."""
        title = content.get("title", "")

        # Full-slide primary background
        self.applier.apply_slide_background(slide, color_key="primary")

        left, top, cw, ch = self._content_rect()

        # Center vertically: start slightly above middle
        title_h = Inches(1.5)
        title_top = top + (ch - title_h) / 2

        self._add_textbox(
            slide, left, title_top, cw, title_h,
            title, text_type="heading", level="h1",
            align=PP_ALIGN.CENTER, color_key="text_on_primary", uppercase=True,
        )

    def _build_two_column(self, slide, content: dict):
        """Two-column layout: heading + two text columns side by side."""
        left_content = content.get("left", {})
        right_content = content.get("right", {})

        left_pad, top_pad, cw, ch = self._content_rect()
        gap = Inches(self.theme.spacing.content_gap / 72)

        col_gap = Inches(0.4)
        col_w = (cw - col_gap) / 2

        current_top = top_pad

        # Slide-level title
        slide_title = content.get("title", "")
        if slide_title:
            title_h = Inches(0.8)
            self._add_textbox(
                slide, left_pad, current_top, cw, title_h,
                slide_title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        col_h = ch - (current_top - top_pad)

        # Left column
        self._build_column_box(
            slide, left_pad, current_top, col_w, col_h, left_content,
        )
        # Right column
        self._build_column_box(
            slide, left_pad + col_w + col_gap, current_top, col_w, col_h, right_content,
        )

    def _build_column_box(self, slide, left, top, width, height, col_content: dict):
        """Build a single column textbox with title, body, and bullets."""
        if not isinstance(col_content, dict):
            col_content = {"body": str(col_content)}

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True

        title = col_content.get("title", "")
        body = col_content.get("body", "")
        bullets = col_content.get("bullets", [])

        paragraphs = []

        # Title
        if title:
            paragraphs.append((title, "heading", "h3", False))

        # Body text (treat as sub-heading if short and no bullets)
        if body:
            paragraphs.append((body, "heading" if not bullets else "body", "h3" if not bullets else None, False))

        # Bullets
        for bullet in bullets:
            paragraphs.append((bullet, "body", None, True))

        if not paragraphs:
            return

        for i, (text, text_type, level, is_bullet) in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.space_after = Pt(10)
            run = p.add_run()
            run.text = text
            self.applier.apply_text_style(run, text_type=text_type, level=level)

            if is_bullet:
                p.level = 0

    def _build_image_left(self, slide, content: dict):
        """Left 40% image, right 60% text content."""
        images = content.get("images", [])
        title = content.get("title", "")
        bullets = content.get("bullets", [])
        body = content.get("body", "")

        left_pad, top_pad, cw, ch = self._content_rect()
        gap = Inches(self.theme.spacing.content_gap / 72)

        img_w = cw * 0.40
        text_w = cw * 0.55
        img_left = left_pad
        text_left = left_pad + img_w + Inches(0.3)

        # Place image on left
        if images:
            img_path = self._resolve_image_path(images[0]["path"])
            if img_path and img_path.exists():
                try:
                    slide.shapes.add_picture(
                        str(img_path), img_left, top_pad,
                        width=img_w, height=ch,
                    )
                except Exception:
                    pass  # fallback: no image

        # Text on right
        current_top = top_pad

        if title:
            title_h = Inches(0.8)
            self._add_textbox(
                slide, text_left, current_top, text_w, title_h,
                title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        content_h = ch - (current_top - top_pad)
        if content_h > Inches(0.3) and (bullets or body):
            box = slide.shapes.add_textbox(text_left, current_top, text_w, content_h)
            tf = box.text_frame
            tf.word_wrap = True

            if bullets:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.space_after = Pt(12)
                    run = p.add_run()
                    run.text = bullet
                    self.applier.apply_text_style(run, text_type="body")
            elif body:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = body
                self.applier.apply_text_style(run, text_type="body")

    def _build_image_full(self, slide, content: dict):
        """Full-bleed image with optional text overlay at bottom."""
        images = content.get("images", [])
        title = content.get("title", "")
        body = content.get("body", "")

        # Try to fill the slide with the image
        if images:
            img_path = self._resolve_image_path(images[0]["path"])
            if img_path and img_path.exists():
                try:
                    pic = slide.shapes.add_picture(
                        str(img_path), Inches(0), Inches(0),
                        width=self._width, height=self._height,
                    )
                    # Send to back so text can overlay
                    spTree = slide.shapes._spTree
                    sp = pic._element
                    spTree.remove(sp)
                    spTree.insert(2, sp)
                except Exception:
                    pass

        # Overlay text at bottom with semi-transparent background bar
        if title or body:
            overlay_h = Inches(1.2)
            overlay_top = self._height - overlay_h - Inches(0.2)
            overlay_w = self._width - Inches(0.4)
            overlay_left = Inches(0.2)

            # Semi-transparent bar (use primary color)
            self.applier.add_shape_with_fill(
                slide, MSO_SHAPE.RECTANGLE,
                overlay_left, overlay_top, overlay_w, overlay_h,
                color_key="primary",
            )

            text = title or body
            self._add_textbox(
                slide, overlay_left + Inches(0.2), overlay_top + Inches(0.2),
                overlay_w - Inches(0.4), overlay_h - Inches(0.4),
                text, text_type="heading", level="h3",
                align=PP_ALIGN.LEFT, color_key="text_on_primary",
            )

    def _build_image_gallery(self, slide, content: dict):
        """Grid layout for multiple images."""
        gallery = content.get("gallery", [])
        title = content.get("title", "")

        left_pad, top_pad, cw, ch = self._content_rect()
        gap = Inches(0.2)

        current_top = top_pad

        # Optional title at top
        if title:
            title_h = Inches(0.8)
            self._add_textbox(
                slide, left_pad, current_top, cw, title_h,
                title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        num_images = len(gallery)
        if num_images == 0:
            return

        # Determine grid dimensions
        if num_images == 2:
            cols = 2
        elif num_images == 3:
            cols = 3
        elif num_images == 4:
            cols = 2
        else:
            cols = 3

        rows = (num_images + cols - 1) // cols
        available_h = ch - (current_top - top_pad)
        cell_w = (cw - gap * (cols - 1)) / cols
        cell_h = (available_h - gap * (rows - 1)) / rows

        for i, item in enumerate(gallery):
            img_path = self._resolve_image_path(item.get("path", ""))
            if not img_path or not img_path.exists():
                continue

            row = i // cols
            col = i % cols
            img_left = left_pad + col * (cell_w + gap)
            img_top = current_top + row * (cell_h + gap)

            try:
                slide.shapes.add_picture(
                    str(img_path), img_left, img_top,
                    width=cell_w, height=cell_h,
                )
            except Exception:
                pass

    def _build_table_slide(self, slide, content: dict):
        """Render a table using shapes."""
        table_data = content.get("table", {})
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        title = content.get("title", "")

        left_pad, top_pad, cw, ch = self._content_rect()
        gap = Inches(self.theme.spacing.content_gap / 72)

        current_top = top_pad

        # Title above table
        if title:
            title_h = Inches(0.8)
            self._add_textbox(
                slide, left_pad, current_top, cw, title_h,
                title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        if not headers and not rows:
            return

        num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
        if num_cols == 0:
            return

        col_w = cw / num_cols
        header_h = Inches(0.5)
        row_h = Inches(0.4)

        from pptx.dml.color import RGBColor
        from html_pptx_template.utils.color import hex_to_rgb

        # Header row background
        header_bg = self.theme.colors.primary
        hr, hg, hb = hex_to_rgb(header_bg)

        # Draw header cells
        for col_idx, header_text in enumerate(headers):
            cell_left = left_pad + col_idx * col_w
            shape = self.applier.add_shape_with_fill(
                slide, MSO_SHAPE.RECTANGLE,
                cell_left, current_top, col_w, header_h,
                color_key="primary",
            )
            # Add border
            shape.line.color.rgb = RGBColor(200, 200, 200)
            shape.line.width = Pt(1)

            # Header text (white, centered)
            self._add_textbox(
                slide, cell_left + Inches(0.05), current_top + Inches(0.05),
                col_w - Inches(0.1), header_h - Inches(0.1),
                header_text, text_type="body",
                align=PP_ALIGN.CENTER, color_key="text_on_primary",
            )

        current_top += header_h

        # Data rows
        for row_idx, row_data in enumerate(rows):
            bg_color_key = "background" if row_idx % 2 == 0 else "surface"
            bgr, bgg, bgb = hex_to_rgb(getattr(self.theme.colors, bg_color_key))

            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    break
                cell_left = left_pad + col_idx * col_w
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, cell_left, current_top, col_w, row_h,
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(bgr, bgg, bgb)
                shape.line.color.rgb = RGBColor(200, 200, 200)
                shape.line.width = Pt(1)

                self._add_textbox(
                    slide, cell_left + Inches(0.05), current_top + Inches(0.05),
                    col_w - Inches(0.1), row_h - Inches(0.1),
                    str(cell_text), text_type="body",
                    align=PP_ALIGN.LEFT, color_key="text_primary",
                )

            current_top += row_h

    def _build_chart_slide(self, slide, content: dict):
        """Chart placeholder slide."""
        chart_info = content.get("chart", {})
        chart_type = chart_info.get("type", "Chart") if isinstance(chart_info, dict) else "Chart"
        title = content.get("title", "")

        left_pad, top_pad, cw, ch = self._content_rect()
        gap = Inches(self.theme.spacing.content_gap / 72)

        current_top = top_pad

        # Title above placeholder
        if title:
            title_h = Inches(0.8)
            self._add_textbox(
                slide, left_pad, current_top, cw, title_h,
                title, text_type="heading", level="h2",
                align=PP_ALIGN.LEFT, uppercase=True,
            )
            current_top += title_h + gap

        # Placeholder rectangle in center of remaining content area
        available_h = ch - (current_top - top_pad)
        placeholder_w = cw * 0.8
        placeholder_h = available_h * 0.8
        placeholder_left = left_pad + (cw - placeholder_w) / 2
        placeholder_top = current_top + (available_h - placeholder_h) / 2

        from pptx.dml.color import RGBColor
        from html_pptx_template.utils.color import hex_to_rgb

        # Draw border rectangle with surface fill
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            placeholder_left, placeholder_top, placeholder_w, placeholder_h,
        )
        fill = shape.fill
        fill.solid()
        sr, sg, sb = hex_to_rgb(self.theme.colors.surface)
        fill.fore_color.rgb = RGBColor(sr, sg, sb)

        # Border with accent color
        ar, ag, ab = hex_to_rgb(self.theme.colors.accent)
        shape.line.color.rgb = RGBColor(ar, ag, ab)
        shape.line.width = Inches(0.2)

        # Centered text inside placeholder
        label_text = f"Chart: {chart_type}"
        self._add_textbox(
            slide, placeholder_left, placeholder_top + placeholder_h / 2 - Inches(0.3),
            placeholder_w, Inches(0.6),
            label_text, text_type="heading", level="h3",
            align=PP_ALIGN.CENTER, color_key="text_primary",
        )
